#!/usr/bin/env python3
"""
Build an MVP PowerPoint from the ABU GIF animations.
Run after saving all GIFs to this folder:
    python3 make_pptx.py
"""

import glob
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ── Slide order and labels ────────────────────────────────────────────────────
SLIDES = [
    ("gif1-frontier-growth.gif",  "GIF 1 — A Frontier Growth"),
    ("gif2-a-to-b.gif",           "GIF 2 — A→B Maturation"),
    ("gif3-u-filtering.gif",      "GIF 3 — U Filtering (combined)"),
    ("gif3-u1-snap.gif",          "GIF 3a — U1 Snap"),
    ("gif3-u2-fade.gif",          "GIF 3b — U2 Fade"),
]

BASE = Path(__file__).parent / "gifs"
OUT  = Path(__file__).parent / "ABU-animations-MVP.pptx"

# ── Slide dimensions: 800×600px ≈ 8.33 × 6.25 in (96 dpi) ───────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Title slide ───────────────────────────────────────────────────────────────
title_slide = prs.slides.add_slide(blank_layout)
bg = title_slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

txBox = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "A / B / U  —  3D Animations"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

sub = tf.add_paragraph()
sub.text = "MVP visual overview"
sub.font.size = Pt(18)
sub.font.color.rgb = RGBColor(170, 170, 170)

# ── One slide per GIF ─────────────────────────────────────────────────────────
missing = []

for filename, label in SLIDES:
    gif_path = BASE / filename
    slide = prs.slides.add_slide(blank_layout)

    # Black background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    if gif_path.exists():
        # Centre the GIF (800×600) on a 960×720 slide with room for label
        img_w = Inches(8)
        img_h = Inches(6)
        left  = (SLIDE_W - img_w) / 2
        top   = Inches(0.5)
        slide.shapes.add_picture(str(gif_path), left, top, img_w, img_h)
    else:
        missing.append(filename)
        # Placeholder text
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"[ GIF not found: {filename} ]"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 68, 68)

    # Caption at bottom
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    cp = cap.text_frame.paragraphs[0]
    cp.text = label
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(200, 200, 200)

prs.save(str(OUT))
print(f"Saved: {OUT}")

if missing:
    print(f"\nMissing GIFs (placeholders added):")
    for f in missing:
        print(f"  {f}")
else:
    print("All GIFs found and embedded.")
